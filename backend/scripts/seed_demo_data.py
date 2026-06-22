import random
import sys
from datetime import datetime, timedelta
from io import BytesIO
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import pandas as pd
from pptx import Presentation

from app.auth import generate_token, hash_password
from app.database import Base, SessionLocal, engine
from app.models import Activity, Assignment, AuditLog, Hackathon, InviteLink, Organizer, Participant, PitchReview, Reviewer, Score, Team
from app.services import registration_intelligence as ri
from app.services.pitch_analyzer import analyze_pitch
from app.services.reviewer_assignment import assign_reviewers

CSV_PATH = Path(__file__).resolve().parent.parent.parent / "sample_data" / "test.csv"

REVIEWER_PANEL = [
    {"name": "Dr. Mehta", "role": "judge", "expertise": ["AI/ML", "Data"], "experience_years": 10, "organization": "IIT Delhi", "industry": "AI/ML", "max_load": 30, "email": "mehta@iitdelhi.demo", "linkedin_url": "https://linkedin.com/in/drmehta", "availability_window": "Sat 10am-6pm", "bio": "AI researcher specializing in NLP."},
    {"name": "Mr. Rao", "role": "judge", "expertise": ["Frontend", "Backend"], "experience_years": 6, "organization": "Infosys", "industry": "Web", "max_load": 30, "email": "rao@infosys.demo", "linkedin_url": "https://linkedin.com/in/mrrao", "availability_window": "Sat-Sun all day", "bio": "Full-stack engineer."},
    {"name": "Ms. Iyer", "role": "judge", "expertise": ["Backend", "Cloud/DevOps", "Data"], "experience_years": 7, "organization": "TCS", "industry": "Cloud", "max_load": 30, "email": "iyer@tcs.demo", "linkedin_url": "https://linkedin.com/in/msiyer", "availability_window": "Sat 9am-5pm", "bio": "Cloud infrastructure lead."},
    {"name": "Mr. Khan", "role": "mentor", "expertise": ["Mobile", "Emerging Tech"], "experience_years": 5, "organization": "Wipro", "industry": "Mobile", "max_load": 30, "email": "khan@wipro.demo", "linkedin_url": "https://linkedin.com/in/mrkhan", "availability_window": "Online most evenings", "bio": "Mobile app developer, loves helping beginners."},
    {"name": "Dr. Suresh", "role": "mentor", "expertise": ["Security", "Programming"], "experience_years": 9, "organization": "Tech Mahindra", "industry": "Security", "max_load": 30, "email": "suresh@techmahindra.demo", "linkedin_url": "https://linkedin.com/in/drsuresh", "availability_window": "Sat 12pm-8pm", "bio": "Security researcher and CTF enthusiast."},
    {"name": "Ms. Patel (Harsh Reviewer)", "role": "judge", "expertise": ["AI/ML", "Backend", "Frontend", "Data", "Cloud/DevOps"], "experience_years": 4, "organization": "Accenture", "industry": "General", "max_load": 30, "email": "patel@accenture.demo", "linkedin_url": "https://linkedin.com/in/mspatel", "availability_window": "Sat 10am-4pm", "bio": "Generalist reviewer."},
]

DEMO_ORGANIZER = {"name": "Demo Organizer", "email": "organizer@demo.dev", "password": "demo1234"}


def build_sample_pitch_deck() -> BytesIO:
    prs = Presentation()
    layout = prs.slide_layouts[1]

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Problem Statement"
    slide1.placeholders[1].text = "Hackathon judges spend hours manually reading every submission before scoring."

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Our Solution"
    slide2.placeholders[1].text = "Python, FastAPI, and scikit-learn power an AI-assisted scoring pipeline deployed on AWS with Docker."

    slide3 = prs.slides.add_slide(layout)
    slide3.shapes.title.text = "Tech Stack"
    slide3.placeholders[1].text = "React frontend, PostgreSQL, pytest test suite, CI/CD pipeline. github.com/example/repo"

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def main():
    random.seed(42)
    Base.metadata.drop_all(bind=engine)
    Base.metadata.create_all(bind=engine)
    db = SessionLocal()

    password_hash, salt = hash_password(DEMO_ORGANIZER["password"])
    organizer = Organizer(
        name=DEMO_ORGANIZER["name"], email=DEMO_ORGANIZER["email"],
        password_hash=password_hash, password_salt=salt, auth_token=generate_token(),
        organization_name="Demo University Hackathon Cell",
        contact_phone="+910000000000",
        logo_url="https://example.com/demo-logo.png",
    )
    db.add(organizer)
    db.flush()

    hackathon = Hackathon(
        organizer_id=organizer.id, name="Demo Hackathon 2026",
        description="A sample hackathon seeded for local demo/testing.",
        sponsors=["Acme Corp", "Globex"], themes=["AI/ML", "Climate Tech"],
        allow_judges=True, allow_mentors=True,
        registration_deadline=datetime.utcnow() + timedelta(days=30),
        mode="hybrid", venue="Demo University Auditorium + online",
        min_team_size=2, max_team_size=5,
        eligibility_criteria="Open to all currently enrolled students.",
        code_of_conduct="Be respectful. No plagiarism. Have fun.",
        support_email="support@demohackathon.dev",
    )
    db.add(hackathon)
    db.commit()
    hackathon_id = hackathon.id
    print(f"Seeded organizer ({DEMO_ORGANIZER['email']} / {DEMO_ORGANIZER['password']}) and hackathon #{hackathon_id}.")
    print(f"Organizer auth_token: {organizer.auth_token}")

    judge_link = InviteLink(hackathon_id=hackathon_id, role="judge", token=generate_token())
    db.add(judge_link)
    db.commit()
    print(f"  Invite link (judge): /join/{judge_link.token}")
    print("  Participants and mentors self-register without an invite link — see /api/hackathons/public")

    print(f"Loading {CSV_PATH} ...")
    df = pd.read_csv(CSV_PATH, dtype={"phone_number": str})
    analyzed = ri.analyze(df)
    metrics = ri.evaluate(analyzed)
    print(f"Ground truth metrics: {metrics['accuracy']=} {metrics['macro_f1']=} {metrics['roc_auc_fraud_vs_genuine']=}")

    teams_by_id_str = {}
    for team_id_str, group in analyzed.groupby("team_id"):
        first = group.iloc[0]
        team = Team(
            hackathon_id=hackathon_id,
            team_id_str=team_id_str,
            team_name=str(first["team_name"]),
            college=str(first["college"]),
            project_idea=str(first["project_idea"]),
            skill_categories=ri.team_skill_categories(group),
            team_size=int(first["team_size"]),
            team_health_score=float(first["team_health_score"]),
            project_novelty_score=float(first["project_novelty_score"]),
        )
        db.add(team)
        teams_by_id_str[team_id_str] = team
    db.flush()

    for _, row in analyzed.iterrows():
        db.add(Participant(
            hackathon_id=hackathon_id,
            registration_id=str(row["id"]),
            team_id=teams_by_id_str[row["team_id"]].id,
            name=str(row["name"]),
            email=str(row["email"]),
            phone_number=str(row.get("phone_number", "")),
            college=str(row["college"]),
            skills=str(row["skills"]),
            project_idea=str(row["project_idea"]),
            team_name=str(row["team_name"]),
            final_trust_score=float(row["final_trust_score"]),
            final_risk_level=str(row["final_risk_level"]),
            confidence_score=float(row["confidence_score"]),
            anomaly_score=float(row["anomaly_score"]),
            fraud_ring_id=int(row["fraud_ring_id"]),
            reasons=list(row["reasons"]),
            explanation=str(row["explanation"]),
            ground_truth_label=str(row["ground_truth_label"]) if pd.notna(row.get("ground_truth_label")) else None,
            predicted_label=str(row["predicted_label"]),
            ip_address=str(row["ip_address"]) if pd.notna(row.get("ip_address")) else None,
            github_username=str(row["github_username"]) if pd.notna(row.get("github_username")) else None,
            academic_year=str(row["academic_year"]) if pd.notna(row.get("academic_year")) else None,
            consent_accepted=True,
        ))
    db.commit()
    print(f"Seeded {len(analyzed)} participants across {len(teams_by_id_str)} teams.")

    reviewers = []
    for spec in REVIEWER_PANEL:
        reviewer = Reviewer(hackathon_id=hackathon_id, auth_token=generate_token(), **spec)
        db.add(reviewer)
        reviewers.append(reviewer)
    db.commit()
    print(f"Seeded {len(reviewers)} reviewers.")

    all_teams = db.query(Team).filter(Team.hackathon_id == hackathon_id).all()
    demo_teams = all_teams[:40]
    judges = [r for r in reviewers if r.role == "judge"]
    results = assign_reviewers(demo_teams, judges)
    assignments = []
    for r in results:
        if r["reviewer_id"] is None:
            continue
        assignment = Assignment(
            hackathon_id=hackathon_id,
            team_id=r["team_id"], reviewer_id=r["reviewer_id"],
            match_score=r["match_score"], explanation=r["explanation"],
        )
        db.add(assignment)
        reviewer = next(rv for rv in reviewers if rv.id == r["reviewer_id"])
        reviewer.current_load += 1
        assignments.append(assignment)
    db.commit()
    print(f"Created {len(assignments)} reviewer assignments.")

    harsh_reviewer = next(r for r in reviewers if "Harsh" in r.name)
    for assignment in assignments:
        reviewer = next(r for r in reviewers if r.id == assignment.reviewer_id)
        base = random.gauss(72, 8)
        score_value = max(0, min(100, round(base, 1)))
        db.add(Score(
            hackathon_id=hackathon_id,
            assignment_id=assignment.id, reviewer_id=reviewer.id, team_id=assignment.team_id,
            score=score_value, criteria={"innovation": score_value, "technical": score_value},
            comments="Auto-generated demo score",
        ))
        db.add(AuditLog(
            hackathon_id=hackathon_id,
            entity_type="score", entity_id=assignment.id, action="submit_score",
            actor=reviewer.name, after={"score": score_value}, reasoning="seed data",
        ))

        if reviewer.id != harsh_reviewer.id:
            harsh_score = max(0, min(100, round(random.gauss(50, 6), 1)))
            db.add(Score(
                hackathon_id=hackathon_id,
                assignment_id=None, reviewer_id=harsh_reviewer.id, team_id=assignment.team_id,
                score=harsh_score, criteria={"innovation": harsh_score, "technical": harsh_score},
                comments="Auto-generated demo score (biased reviewer)",
            ))
            db.add(AuditLog(
                hackathon_id=hackathon_id,
                entity_type="score", entity_id=assignment.id, action="submit_score",
                actor=harsh_reviewer.name, after={"score": harsh_score}, reasoning="seed data (biased)",
            ))
    db.commit()
    print("Seeded synthetic scores, including a deliberately biased reviewer:", harsh_reviewer.name)

    now = datetime.utcnow()
    active_teams = demo_teams[:15]
    declining_teams = demo_teams[15:18]
    inactive_teams = demo_teams[18:20]
    untouched_teams = demo_teams[20:22]

    for team in active_teams:
        for hours_ago in [2, 8, 20, 30, 40]:
            db.add(Activity(hackathon_id=hackathon_id, team_id=team.id, activity_type="commit", timestamp=now - timedelta(hours=hours_ago)))

    for team in declining_teams:
        for hours_ago in [85, 80, 75, 70]:
            db.add(Activity(hackathon_id=hackathon_id, team_id=team.id, activity_type="commit", timestamp=now - timedelta(hours=hours_ago)))
        db.add(Activity(hackathon_id=hackathon_id, team_id=team.id, activity_type="commit", timestamp=now - timedelta(hours=10)))

    for team in inactive_teams:
        db.add(Activity(hackathon_id=hackathon_id, team_id=team.id, activity_type="commit", timestamp=now - timedelta(hours=70)))

    db.commit()
    print(
        f"Seeded activity logs: {len(active_teams)} active, {len(declining_teams)} declining, "
        f"{len(inactive_teams)} inactive, {len(untouched_teams)} with no activity at all."
    )

    sample_team = demo_teams[0]
    deck = build_sample_pitch_deck()
    pitch_result = analyze_pitch(deck, project_novelty_score=sample_team.project_novelty_score)
    db.add(PitchReview(hackathon_id=hackathon_id, team_id=sample_team.id, **pitch_result))
    db.commit()
    print(f"Seeded a sample pitch review for team {sample_team.team_name}.")

    pending_participants = [
        Participant(
            hackathon_id=hackathon_id, registration_id=generate_token()[:16], auth_token=generate_token(),
            name="Late Joiner One", email="latejoiner1@demo.dev", team_name="LastMinuteCrew",
            skills="Python,Flask", project_idea="Campus Lost & Found Tracker",
            ip_address="192.0.2.10", approval_status="pending",
            github_username="latejoiner_demo", academic_year="3rd Year", consent_accepted=True,
            reasons=[], explanation="Awaiting full risk analysis — instant checks only: no red flags",
        ),
        Participant(
            hackathon_id=hackathon_id, registration_id=generate_token()[:16], auth_token=generate_token(),
            name="Late Joiner Two", email="latejoiner2@demo.dev", team_name="LastMinuteCrew",
            skills="Python,Flask", project_idea="Campus Lost and Found Tracker App",
            ip_address="192.0.2.10", approval_status="pending",
            github_username="latejoiner_demo", academic_year="3rd Year", consent_accepted=True,
            reasons=[], explanation="Awaiting full risk analysis — instant checks only: no red flags",
        ),
    ]
    for p in pending_participants:
        db.add(p)
    db.commit()
    print(f"Seeded {len(pending_participants)} pending self-registered participants awaiting organizer review.")

    db.close()
    print("Done. Start the API with: uvicorn app.main:app --reload")


if __name__ == "__main__":
    main()
