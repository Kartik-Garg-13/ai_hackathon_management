from dataclasses import dataclass

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.skill_taxonomy import SKILL_TAXONOMY

TESTING_KEYWORDS = ["test", "pytest", "unit test", "testing", "qa", "ci/cd"]
DEPLOYMENT_KEYWORDS = ["deploy", "docker", "kubernetes", "aws", "azure", "gcp", "hosting", "render", "vercel", "netlify"]
PROBLEM_KEYWORDS = ["problem", "challenge", "pain point", "issue we", "the gap"]
EMERGING_TECH_CATEGORIES = {"AI/ML", "Emerging Tech"}


@dataclass
class SlideStats:
    slide_count: int
    image_count: int
    full_text: str
    word_count: int


def extract_slide_stats(pptx_path_or_file) -> SlideStats:
    prs = Presentation(pptx_path_or_file)
    texts = []
    image_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
    full_text = "\n".join(texts)
    return SlideStats(
        slide_count=len(prs.slides),
        image_count=image_count,
        full_text=full_text,
        word_count=len(full_text.split()),
    )


def detect_technical_complexity(text: str) -> tuple[str, int]:
    text_l = text.lower()
    found = {term for term in SKILL_TAXONOMY if term in text_l}
    count = len(found)
    if count >= 6:
        return "High", count
    if count >= 3:
        return "Medium", count
    return "Low", count


def detect_presentation_quality(stats: SlideStats) -> str:
    if stats.slide_count == 0:
        return "Low"
    words_per_slide = stats.word_count / stats.slide_count
    images_per_slide = stats.image_count / stats.slide_count
    if 15 <= words_per_slide <= 80 and images_per_slide >= 0.3:
        return "High"
    if words_per_slide < 8 or words_per_slide > 150:
        return "Low"
    return "Medium"


def detect_issues(text: str) -> list[str]:
    text_l = text.lower()
    issues = []
    if not any(kw in text_l for kw in TESTING_KEYWORDS):
        issues.append("No testing found")
    if not any(kw in text_l for kw in DEPLOYMENT_KEYWORDS):
        issues.append("No deployment strategy mentioned")
    if not any(kw in text_l for kw in PROBLEM_KEYWORDS):
        issues.append("Problem statement unclear")
    if "github.com" not in text_l and "http" not in text_l:
        issues.append("No project link found")
    return issues


def compute_innovation_score(text: str, project_novelty_score: float | None) -> float:
    text_l = text.lower()
    emerging_hits = sum(
        1 for term, category in SKILL_TAXONOMY.items()
        if category in EMERGING_TECH_CATEGORIES and term in text_l
    )
    novelty_component = (project_novelty_score or 50.0) / 10
    emerging_component = min(emerging_hits, 5) / 5 * 10
    return round(0.6 * novelty_component + 0.4 * emerging_component, 1)


def analyze_pitch(pptx_path_or_file, project_novelty_score: float | None = None) -> dict:
    stats = extract_slide_stats(pptx_path_or_file)
    technical_complexity, _ = detect_technical_complexity(stats.full_text)
    presentation_quality = detect_presentation_quality(stats)
    issues = detect_issues(stats.full_text)
    innovation_score = compute_innovation_score(stats.full_text, project_novelty_score)

    return {
        "innovation_score": innovation_score,
        "technical_complexity": technical_complexity,
        "presentation_quality": presentation_quality,
        "issues": issues,
        "slide_count": stats.slide_count,
    }
