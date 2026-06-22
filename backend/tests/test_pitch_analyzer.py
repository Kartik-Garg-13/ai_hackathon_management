from io import BytesIO

from pptx import Presentation

from app.services.pitch_analyzer import analyze_pitch


def build_deck(slide_texts: list[str]) -> BytesIO:
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for title, body in slide_texts:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def test_strong_deck_scores_high_complexity_and_no_issues():
    deck = build_deck([
        ("Problem", "Judges face a real challenge reviewing hundreds of projects manually."),
        ("Solution", "We use Python, TensorFlow, BERT, and FastAPI deployed on AWS with Docker."),
        ("Quality", "Covered with pytest unit tests and a CI/CD pipeline. github.com/example/repo"),
    ])
    result = analyze_pitch(deck, project_novelty_score=80)
    assert result["technical_complexity"] == "High"
    assert result["issues"] == []
    assert result["slide_count"] == 3


def test_weak_deck_flags_all_issues():
    deck = build_deck([("My Project", "Cool app idea.")])
    result = analyze_pitch(deck, project_novelty_score=20)
    assert result["technical_complexity"] == "Low"
    assert "No testing found" in result["issues"]
    assert "No deployment strategy mentioned" in result["issues"]
    assert "Problem statement unclear" in result["issues"]
    assert "No project link found" in result["issues"]


def test_innovation_score_responds_to_novelty():
    deck = build_deck([("Title", "Some text about our project.")])
    high_novelty = analyze_pitch(deck, project_novelty_score=95)
    deck.seek(0)
    low_novelty = analyze_pitch(deck, project_novelty_score=5)
    assert high_novelty["innovation_score"] > low_novelty["innovation_score"]
